from pptx import Presentation #lets us create presentation
from pptx.util import Inches, Pt #lets us get pptx fonts, measurements etc
import time #lets us pause for a certain time
from pydrive.drive import GoogleDrive #Lets us upload to google drive
from pydrive.auth import GoogleAuth # Authentication for google account
import urllib # lets us download image from ur
import urllib2
import os # lets us execute console commands
from datetime import datetime  
from selenium import webdriver #lets us search the web
from selenium.webdriver.common.by import By #lets us search for elements
from selenium.webdriver.common.keys import Keys #lets us enter keys into webdriver
from selenium.webdriver.remote.webelement import WebElement 


############################################################################################################################################################
############################################################################################################################################################
############################################################################################################################################################
#
# HEYAW METEKE
# ALGEBRA 2
# MS HELFT
# E BLOCK
#
############################################################################################################################################################
############################################################################################################################################################
############################################################################################################################################################


class loan (object):
        simple = {'rate': 0.04}
        compound = {'rate': 1.0581}
	collegeTuition = 0
        def __init__(self, tuition):
		print 'Calculating Loans...',
		self.collegeTuition = tuition
		self.find_interest()

        def find_interest (self):

		##############################################################################
		############## FINDS INTEREST FOR BOTH COMPOUND AND SIMPLE ###################
		##############################################################################

                self.simple['equation'] = 'Interest = ' + str(self.collegeTuition) + '(' + str(self.simple['rate']) + '(year)'
		self.compound['equation'] = 'Interest = ' + str(self.collegeTuition) + '(' + str(self.compound['rate']) + ')^' + 'year)'

            
                for year in range(0, 6): #LOOP FINDS INTEREST FOR 6 YEARS 
                        if year >= 0 and year <= 5:

				# FINDS simple interest
                                full = self.collegeTuition * self.simple['rate'] * year
                                self.simple[year] = str(full)

				#FINDS compound interest
				full = self.collegeTuition * ( (self.compound['rate']) ** year)
				self.compound[year] = str(full)
			else:
				pass

		print '\tOK'

class powerpoint (object):
	
	def __init__(self):
		print 'Creating Slides...',
		###########################################################################
		#	 FIRST TIME RUN TO CREATE PRESENTATION AND MAKE DIRECTORY	  # 
		###########################################################################

		
		title_page_layout = prs.slide_layouts[0] #Creates
		title_slide = prs.slides.add_slide(title_page_layout)
		title_page_title = title_slide.shapes.title
		title_page_subtitle = title_slide.placeholders[1]		
		title_page_title.text = 'College Loan Project'
		title_page_subtitle.text = 'Programmed by Heyaw Meteke'
	
		
		###########################################################################
		#	EACH FUNCTION CREATES PART OF POWER POINT PRESENTATION		  #
		###########################################################################

		self.add_project_page()
		self.add_image_page()
		self.add_info_page()
		self.add_loan_page()
		self.add_cost_page()
		self.export()

	
	def export (self): 
		
		######################################################
		########## EXPORTS FILE TO GOOGLE DRIVE ##############
		#####################################################		

		gauth = GoogleAuth()
		gauth.LoadCredentialsFile("creds.txt")

		if gauth.credentials is None:
			gauth.LocalWebserverAuth()
		
		elif gauth.access_token_expired:
			gauth.Refresh()

		else:
			gauth.Authorize()

		gauth.SaveCredentialsFile("creds.txt")
		gdrive = GoogleDrive(gauth)
		
		name = start.collegeName
		folder_id = '0B22isVxRZHa5aG1nMWpkQzdZRzQ'
		prs.save('slides.pptx')
			
		try:
			print 'Uploading File...',
			presentation_file = gdrive.CreateFile({'title':name, 'parents': [{'kind': 'drive#fileLink', 'id': folder_id}]})
			presentation_file.SetContentFile('slides.pptx')
			presentation_file.Upload()
			print '\tOK'
		
			print 'Cleaning Up Files...',
			os.system('rm slides.pptx')
			os.system('rm image.png')
			print '\tOK'
			print 'Copy and paste the following link to go to folder'
			print 'LINK: https://goo.gl/akeZ3q'	
		
		except:
			print 'ERROR! FILE COULD NOT UPLOAD!!!'
			print 'File is in program directory.'


	def add_image_page(self):
		
		#########################################################
		############# CREATES AND SETS IMAGE PAGE ###############
		########################################################

		image_page_layout = prs.slide_layouts[6]
		image_slide =prs.slides.add_slide(image_page_layout)
		image_path = 'image.png'
	
		top = Inches(1)
		height = Inches(5.5)
		left = Inches(2)
		
		for tries in range (0,3):

			try:
				img = image_slide.shapes.add_picture(image_path, left, top, height=height)


			except:
				print 'Image is corrupted, adding placeholder.'
				img = 'placeholder.png'	
				image = image_slide,shapes.add_pictures(image_path, left, top, height=height)	
			else:
				break

	def add_project_page (self):
		slide_layout = prs.slide_layouts[6]
		project_slide = prs.slides.add_slide(slide_layout)
		title = project_slide.shapes.title
#		title.text = 'Loans'
		
		left = top = Inches(1)
		width = height = Inches(8)
		txBox = project_slide.shapes.add_textbox(left, top, width, height)
		tf = txBox.text_frame
		
		temp_text = ''
		with open('explain.txt') as f:
			temp_text = f.readlines()
		tf.text = unicode(''.join(temp_text), 'utf-8')

	def add_info_page (self):

		#######################################################################
		################## CREATES AND SETS INFORMATION PAGE ##################
		######################################################################		
				
		name = start.collegeName
		location = start.collegeLocation
		desc = start.collegeDesc
		tuition = start.collegeTuition

		info_page_layout = prs.slide_layouts[1]
		info_slide = prs.slides.add_slide(info_page_layout)
		modules = info_slide.shapes
	
		title_info_page = modules.title
		body_info_page = modules.placeholders[1]
		
		title_info_page.text = name
		
		textbox = body_info_page.text_frame
		textbox.text = 'Location: ' + location
	
		t = textbox.add_paragraph()
		t.text = 'Tuition (Including fees/room and board): $' + str(tuition)
		t.font.size = Pt(20)
	
		d = textbox.add_paragraph()
		d.text = desc
		d.font.size = Pt(15)
		
	def add_loan_page (self):
			
		
                slide_layout = prs.slide_layouts[6]
                loan_slide = prs.slides.add_slide(slide_layout)
                
		left = top = Inches(1)
                width = height = Inches(8)
                txBox = loan_slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame

                temp_text = ''
                with open('loaninfo.txt') as f:
                        temp_text = f.readlines()
                temp_text = unicode(''.join(temp_text), 'utf-8')
		
		t = str(start.collegeTuition)
		vocab = {1: start.collegeName, 2: t, 3: t, 4: t , 5: t, 6: t}
		counter = 1
		for char in temp_text:
			if counter >=1 and counter <= 6:
				index = temp_text.find('#')
				temp_text = temp_text[:index-1] + vocab[counter] + temp_text[index+1:]		
				counter +=1
		tf.text = temp_text
	def add_cost_page (self):

		################################################################################################
		################### CREATES AND SETS TABLES FOR COMPOUND AND SIMPLE INTEREST####################
		################################################################################################

		simple_loan = loan.simple
		compound_loan = loan.compound

		tuition = start.collegeTuition

		#CREATES SLIDE FOR BOTH COMPOUND AND SIMPLE
		table_page_layout = prs.slide_layouts[5]
		
		simple_slide = prs.slides.add_slide(table_page_layout)
		compound_slide = prs.slides.add_slide(table_page_layout)

		left = top = height = Inches(1)
		width = Inches(2)
		
	
		modules_simple = simple_slide.shapes
		modules_compound = compound_slide.shapes


		# ADDS TITLE FOR BOTH COMPOUND AND SIMPLE SLIDES
		modules_simple.title.text = 'Simple Interest'

		modules_compound.title.text = 'Compound Interest'
	#csubtitle.text = compound_loan['equation']
		
		rows = 7
		cols = 3
		left = top = Inches(2.0)
		width = Inches(6.0)
		height = Inches(0.8)
		
		# CREATES BOTH COMPOUND AND SIMPLE TABLE
		simple_table = modules_simple.add_table(rows, cols, left, top, width, height).table
		compound_table = modules_compound.add_table(rows, cols, left, top, width, height).table		
		
		# SETS HEADER FOR BOTH COMPOUND AND SIMPLE TABLE
		compound_table.cell(0, 0).text = simple_table.cell(0, 0).text = 'Year' 
		compound_table.cell(0, 1).text = simple_table.cell(0, 1).text = 'Interest'
		compound_table.cell(0, 2).text = simple_table.cell(0, 2).text = 'Balance'
		
		counter = 1
		
		tuition = int(tuition)

		for year in range(0,6):
	
			if year >= 0 and year <= 5 :
				# ADDS VALUES FOR SIMPLE TABLE
				simple_table.cell(counter, 0).text = str(year)
				simple_table.cell(counter, 1).text = str(simple_loan[year]) 
				simple_table.cell(counter, 2).text = str(round(float(simple_loan[year])+tuition, 2))
			
				# ADDS VALUES FOR COMPOUND TABLE
				compound_table.cell(counter, 0).text = str(year)

				if isinstance(compound_loan[year], basestring): # checks if an interest is a string, if yes it changes to float
					compound_loan[year] = float(compound_loan[year])


				compound_table.cell(counter, 1).text = str(compound_loan[year]-tuition)
				compound_table.cell(counter, 2).text = str(round(compound_loan[year], 2))
				
				counter+=1				

		print '\tOK'

class college_scrapper (object):
	collegeName = ''
	collegeImageDir = ''
	collegeLocation = ''
	collegeTuition = 0
	collegeDesc = ''' '''
	collegDir = ''
	loan_table = {}
	url = ''
	name_search = ''
	def __init__ (self, name_search):
		os.system('clear')
		self.name_search = name_search

		print '\t\t\tCollege Loan Algebra Project'
		print '\t\t\tProgrammed by Heyaw Meteke'	
		
		driver.get('http://www.google.com')
		driver.set_window_size(1440, 900)
	
	
		self.college_search()
		self.get_image()

        def get_image (self):
		image_url = ''
                driver.get('https://www.yandex.com/images/')
                search_image = driver.find_element_by_name('text')
                search_image.clear()
		search_image.send_keys('logo ' + self.collegeName)
		search_image.send_keys(Keys.RETURN)

		
		try:
			time.sleep(4)
			print 'Searching for Image...',
			temp_image_url = driver.find_element_by_css_selector("*[class^='serp-item__link']").get_attribute('href')
			driver.get(temp_image_url)
		except:
			time.sleep(3)
			temp_image_url = driver.find_element_by_css_selector("*[class^='serp-item__link']").get_attribute('href')
			driver.get(temp_image_url)
		try:
			#Makes sure to close pop up if its open
			driver.find_element_by_css_selector("*[class^='popup__content']").click()
		except:
			#if no pop up, no worry - just to catch error
			pass
		try:
	
			image_url = driver.find_element_by_css_selector("*[class^='button2 button2_theme_action button2_size_m button2_type_link button2_pin_brick-clear button2_width_max sizes__download i-bem button2_js_inited']").get_attribute('href')

		except:
			print "First method didnt work, going to second one"
			image_url = driver.find_element_by_xpath("/html/body/div[5]/div[3]/div[2]/div[1]/div[2]/div[1]/a").get_attribute('href')	
		print '\tOK'
		
		print 'Downloading Image...',
		
		self.collegeImageDir =  'image.png'

		urllib.urlretrieve(image_url, 'image.png')
		print '\tOK'

	def college_search (self):

		college_input = self.name_search + ' college data'
		print 'Loading Program...',

		time.sleep(2)
		search_college = driver.find_element_by_name('q')
		search_college.send_keys(college_input)
		search_college.send_keys(Keys.RETURN)
		time.sleep(3)
		
		driver.find_element_by_partial_link_text('CollegeData').click()

		time.sleep(5)		

		self.collegeName = driver.find_element_by_xpath("//*[@id='collprofile']/div[6]/div[4]/div[2]/div[1]/h1").text	# NAME OF COLLEGE

		self.collegeLocation = driver.find_element_by_css_selector("*[class^='citystate']").text
		
		self.collegeTuition = driver.find_element_by_xpath("//*[@id='section1']/table/tbody/tr[1]/td").text # TUITION COST FOR COLLEGE

		try:
			self.collegeDesc = driver.find_element_by_xpath("//*[@id='cont_overview']/p").text #BRIEF DESCRIPTION OF COLLEGE
			print '\tOK'
			
		except: #If it cant find desc on college data - goes to college board instead
			print '\tERROR!'
			print '==========================ERROR==========================================='
			print 'Usual source does not have all information.  Checking other source instead'
			print '==========================================================================\n'
			driver.get('www.google.com')
			time.sleep(2)
			college_input = self.name_search + ' college board'
			driver.save_screenshot('SCREEN.png')
			self.collegeDesc = driver.find_element_by_css_selector("*[class^='spotlighttext']").text
			# search_college.send_keys(college_input)
	                #search_college.send_keys(Keys.RETURN)
			#time.sleep(3)
			#driver.find_element_by_partial_link_text('...').click()
			#time.sleep(3)
			#self.collegeDesc = driver.find_element_by_id("cpProfile_ataglance_collegeDescription_html").text # VERY SHORT DESCRIPTION
	

		#####################################################################################
		################# SEPERATE DATA AND TO SEND TUITION TO LOAN OBJECT ##################
		#####################################################################################	

		self.collegeDir = self.collegeName.replace(' ', "_")

		self.collegeTuition = self.collegeTuition.split('$')

		tuition = ''

		tuitionTemp = self.collegeTuition[1]

		tuitionTemp = tuitionTemp.replace(' ', "")

		for char in tuitionTemp:
			try:
				works = int(char)
				tuition+= char
			except:
				pass
		self.collegeTuition = int(tuition)


###########################################################################################
################################## CREATES INSTANCES ######################################
###########################################################################################
name_search = raw_input('College: ')							###
											###
driver = webdriver.PhantomJS() # FOR FINAL VERSION					###
											###
start = college_scrapper(name_search) # STARTS WEB SCRAPPING				###				
											###
loan = loan(start.collegeTuition) # CALCULATES LOAN FOR BOTH COMPOUND AND SIMPLE	###
											###
prs = Presentation() # CREATES SLIDE INSTANCE						###
											###
p = powerpoint() # ADDS INFORMATION TO SLIDE						###
###########################################################################################
##################################### END #################################################
###########################################################################################
